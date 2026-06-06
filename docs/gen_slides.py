"""
gen_slides.py - 8-slide presentation based on Template-Presentation.pptx.
Run: python docs/gen_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

ROOT     = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS     = os.path.join(ROOT, "docs", "figures")
TEMPLATE = os.path.join(ROOT, "docs", "Template for thesis", "Template-Presentation.pptx")
OUT      = os.path.join(ROOT, "docs", "Presentation_AI_Rubiks_Cube.pptx")

# Color palette (light theme)
NAVY   = RGBColor(0x0D, 0x1B, 0x4B)
BLUE   = RGBColor(0x1A, 0x3A, 0x8F)
STEEL  = RGBColor(0x3A, 0x5A, 0xAA)
BODY   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY   = RGBColor(0x55, 0x55, 0x77)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xFF, 0xFF, 0xFF)
STRIP  = RGBColor(0xE8, 0xEC, 0xF8)
GREEN  = RGBColor(0x1A, 0x7A, 0x40)
RED    = RGBColor(0xAA, 0x22, 0x22)
HDRBG  = RGBColor(0x0D, 0x1B, 0x4B)

W = Inches(13.33)
H = Inches(7.50)

BLANK_LAYOUT = 55   # confirmed from template inspection


# ---------------------------------------------------------------------------
# Template helpers
# ---------------------------------------------------------------------------

def _remove_slides_after_first(prs):
    """Remove slides 1..N keeping only slide 0 (the front page)."""
    sldIdLst = prs.slides._sldIdLst
    while len(sldIdLst) > 1:
        sId = sldIdLst[1]
        rId = sId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sId)


def _update_front_page_date(prs):
    """Replace 'April, 2026' / 'April 2026' with 'May 2026' on slide 0."""
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "April" in run.text:
                    run.text = run.text.replace("April, 2026", "May 2026")
                    run.text = run.text.replace("April 2026", "May 2026")


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txb(slide, l, t, w, h, text, size=14, bold=False, italic=False,
        color=BODY, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def title_bar(slide, text, sub=None):
    """Navy top bar with white title and optional subtitle."""
    add_rect(slide, 0, 0, W, Inches(0.82), NAVY)
    txb(slide, Inches(0.4), Inches(0.06), Inches(12.5), Inches(0.55),
        text, size=26, bold=True, color=WHITE)
    if sub:
        txb(slide, Inches(0.42), Inches(0.60), Inches(12.5), Inches(0.28),
            sub, size=12, italic=True, color=RGBColor(0xBB, 0xBB, 0xDD))


def section_label(slide, text, l, t, w=Inches(4.5)):
    add_rect(slide, l, t, w, Inches(0.30), NAVY)
    txb(slide, l + Inches(0.08), t + Inches(0.02), w, Inches(0.28),
        text, size=11, bold=True, color=WHITE)


def formula_box(slide, l, t, w, h, lines, size=14):
    add_rect(slide, l, t, w, h, STRIP, RGBColor(0xAA, 0xBB, 0xDD))
    box = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.08),
                                   w - Inches(0.3), h - Inches(0.16))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY


def add_pic(slide, fname, l, t, w, h=None):
    path = os.path.join(FIGS, fname)
    if not os.path.exists(path):
        return None
    return slide.shapes.add_picture(path, l, t, w, h) if h else \
           slide.shapes.add_picture(path, l, t, w)


def tbl_create(slide, data, l, t, w, row_h=Inches(0.38)):
    rows = len(data)
    cols = len(data[0])
    tbl  = slide.shapes.add_table(rows, cols, l, t, w, row_h * rows).table
    tbl.first_row = True
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            tf   = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run  = tf.paragraphs[0].add_run()
            run.text = val
            run.font.name = "Calibri"
            run.font.size = Pt(13)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = HDRBG
                run.font.color.rgb = WHITE
                run.font.bold = True
            else:
                cell.fill.fore_color.rgb = (STRIP if r % 2 == 0 else BG)
                if "100%" in val:  run.font.color.rgb = GREEN
                elif "0%" in val:  run.font.color.rgb = RED
                else:              run.font.color.rgb = BODY
    return tbl


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Content slides
# ---------------------------------------------------------------------------

def slide_system_overview(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "System Overview",
              "Four layers that work together from raw state to animated web solve")

    components = [
        ("Cube Engine",       "core/",
         "cp, co, ep, eo arrays\n18 move tables\nLegality checks"),
        ("Expert Solver",     "solvers/kociemba_solver.py",
         "Bidirectional BFS\nGenerates training labels\nBaseline for UI"),
        ("Neural Solver",     "solvers/ai_solver.py",
         "MLP  324->256->256->128->18\nCurriculum imitation learning\nGreedy + beam search"),
        ("Web Visualisation", "visualization/",
         "Python HTTP server\nThree.js 3D frontend\nREST API + animation queue"),
    ]
    for i, (name, path, detail) in enumerate(components):
        lx = Inches(0.3 + i * 3.26)
        add_rect(s, lx, Inches(1.0), Inches(3.1), Inches(5.8),
                 STRIP, RGBColor(0x99, 0xAA, 0xCC))
        add_rect(s, lx, Inches(1.0), Inches(3.1), Inches(0.45), BLUE)
        txb(s, lx + Inches(0.08), Inches(1.03), Inches(2.95), Inches(0.4),
            name, size=15, bold=True, color=WHITE)
        txb(s, lx + Inches(0.1), Inches(1.55), Inches(2.95), Inches(0.3),
            path, size=10, italic=True, color=STEEL)
        txb(s, lx + Inches(0.1), Inches(1.92), Inches(2.9), Inches(4.5),
            detail, size=13, color=BODY)

    for i in range(3):
        lx = Inches(3.38 + i * 3.26)
        add_rect(s, lx, Inches(3.75), Inches(0.22), Inches(0.22), NAVY)
        txb(s, lx, Inches(3.72), Inches(0.22), Inches(0.28),
            "▶", size=12, color=WHITE, align=PP_ALIGN.CENTER)

    set_notes(s,
        "This slide shows the four main components.\n\n"
        "The Cube Engine is the mathematical heart - it stores the cube as four integer arrays "
        "(cp=corner permutation, co=corner orientation, ep=edge permutation, eo=edge orientation) "
        "and applies moves using precomputed permutation tables. Nothing graphical lives here.\n\n"
        "The Expert Solver is a bidirectional BFS that finds optimal solutions for shallow scrambles. "
        "It serves two roles: as the training oracle (generates labels) and as the reliable baseline in the UI.\n\n"
        "The Neural Solver is a 3-hidden-layer MLP trained via curriculum imitation learning. "
        "It has two inference modes: greedy (picks the single best move) and beam search (keeps k candidate paths).\n\n"
        "The Web Visualisation is a Python HTTP server plus a Three.js single-page app. "
        "The server exposes a REST API; the frontend renders 26 independent 3D cubie meshes and animates them.")


def slide_cube_engine(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Cube Engine - State Representation",
              "The cube is stored as four integer arrays - no 3D geometry, no colours")

    section_label(s, "State arrays", Inches(0.3), Inches(1.0))
    arr_data = [
        ["Array", "Size", "Values", "Meaning"],
        ["cp",  "8",  "0 - 7",  "Which corner piece occupies each corner slot"],
        ["co",  "8",  "0, 1, 2","Twist of each corner (0=upright, 1=CW, 2=CCW)"],
        ["ep",  "12", "0 - 11", "Which edge piece occupies each edge slot"],
        ["eo",  "12", "0, 1",   "Whether each edge is flipped (0=normal, 1=flipped)"],
    ]
    tbl_create(s, arr_data, Inches(0.3), Inches(1.40), Inches(7.6), row_h=Inches(0.42))

    add_pic(s, "fig_03_encoding.png", Inches(8.1), Inches(1.0), Inches(4.9))

    txb(s, Inches(0.3), Inches(3.70), Inches(7.6), Inches(0.55),
        "Solved state: cp = [0,1,2,3,4,5,6,7],  co = [0]*8,  ep = [0..11],  eo = [0]*12",
        size=12, italic=True, color=GRAY)

    set_notes(s,
        "The cube engine never stores colours - it stores piece identities.\n\n"
        "cp (corner permutation): cp[3] = 5 means the corner piece originally called piece-5 "
        "is currently sitting in slot 3.\n\n"
        "co (corner orientation): each corner has three-fold rotational symmetry around its "
        "body diagonal. Orientation 0 = reference facelet (the U/D sticker) is on U or D face. "
        "Orientation 1 = rotated 120 degrees clockwise. Orientation 2 = 120 degrees counter-clockwise.\n\n"
        "ep/eo: same idea for edges (12 pieces, binary orientation - flipped or not).\n\n"
        "The three invariants are the group-theoretic constraints that make the cube a genuine group. "
        "They also define which states are reachable: there are exactly 43,252,003,274,489,856,000 "
        "legal configurations, derived by dividing the unconstrained count (8! x 3^8 x 12! x 2^12) by 12.")


def slide_move_tables(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Move Tables - 18 Legal Moves",
              "Each face has 3 variants: CW (90 deg), CCW (90 deg prime), and 180 deg double")

    section_label(s, "The 18-move set (Half-Turn Metric)",
                  Inches(0.3), Inches(1.0), w=Inches(4.2))
    moves_data = [
        ["Face",  "CW", "CCW (prime)", "180 deg"],
        ["Up",    "U",  "U'",          "U2"],
        ["Down",  "D",  "D'",          "D2"],
        ["Right", "R",  "R'",          "R2"],
        ["Left",  "L",  "L'",          "L2"],
        ["Front", "F",  "F'",          "F2"],
        ["Back",  "B",  "B'",          "B2"],
    ]
    tbl_create(s, moves_data, Inches(0.3), Inches(1.40), Inches(5.2), row_h=Inches(0.40))

    section_label(s, "Move composition (group product)",
                  Inches(5.8), Inches(1.0), w=Inches(4.0))
    formula_box(s, Inches(5.8), Inches(1.40), Inches(7.1), Inches(2.65), [
        "# Applying move B then move A:",
        "cp_out[i] = cp_A[ cp_B[i] ]",
        "co_out[i] = (co_A[cp_B[i]] + co_B[i])  mod 3",
        "",
        "ep_out[i] = ep_A[ ep_B[i] ]",
        "eo_out[i] = (eo_A[ep_B[i]] + eo_B[i])  mod 2",
    ], size=12)

    section_label(s, "Prime and double moves derived algebraically",
                  Inches(5.8), Inches(4.25), w=Inches(5.0))
    formula_box(s, Inches(5.8), Inches(4.65), Inches(7.1), Inches(1.35), [
        "R2  = compose(R, R)          # apply R twice",
        "R'  = compose(R2, R)         # apply R three times = inverse",
    ], size=12)

    section_label(s, "Branching factor", Inches(0.3), Inches(4.25), w=Inches(2.5))
    txb(s, Inches(0.3), Inches(4.65), Inches(5.1), Inches(1.1),
        "18 moves at each step, but consecutive same-face moves are pruned.\n"
        "Effective branching factor: ~15. BFS nodes at depth d: ~15^d.",
        size=13, color=BODY)

    set_notes(s,
        "Each of the 18 moves is stored as a MoveDefinition: two permutation tuples (corner_perm and "
        "edge_perm) and two orientation-delta tuples (corner_orient and edge_orient).\n\n"
        "The composition formulas on this slide are the actual code logic in cube.py apply_move():\n"
        "  new_cp[i] = self.cp[ move_def.corner_perm[i] ]\n"
        "  new_co[i] = (self.co[ move_def.corner_perm[i] ] + move_def.corner_orient[i]) % 3\n\n"
        "Prime moves are derived by composing the base move three times (M3 = M^-1 for 90-deg moves). "
        "Double moves are M composed with itself. This keeps the move table algebraically consistent.\n\n"
        "Move pruning: never follow a move with its inverse on the same face (e.g. R then R'), "
        "and for commuting opposite faces (R and L), enforce a fixed order. This reduces the "
        "effective branching factor from 18 to ~15.")


def slide_encoding(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "State Encoding for the Neural Network",
              "The 54-character facelet string is converted to a 324-dimensional one-hot vector")

    section_label(s, "One-hot encoding scheme", Inches(0.3), Inches(1.0), w=Inches(3.5))
    formula_box(s, Inches(0.3), Inches(1.40), Inches(5.8), Inches(2.65), [
        "# For each of the 54 facelets:",
        "xᵢ  ∈  {0.0, 1.0}^6",
        "",
        "# Exactly one element is 1.0:",
        "Σⱼ xᵢ[ⱼ]  =  1   for all i",
        "",
        "# Full input vector:",
        "v  =  concat(x₀, x₁, ..., x₅₃)   ∈  ℝ^324",
    ], size=13)

    section_label(s, "Dimension breakdown", Inches(0.3), Inches(4.25), w=Inches(3.0))
    dim_data = [
        ["Quantity",             "Count"],
        ["Facelets",             "54"],
        ["Colours per facelet",  "6"],
        ["Input vector length",  "54 x 6 = 324"],
        ["Non-zero per state",   "54 (exactly one per facelet)"],
    ]
    tbl_create(s, dim_data, Inches(0.3), Inches(4.65), Inches(5.8), row_h=Inches(0.38))

    add_pic(s, "fig_03_encoding.png", Inches(6.4), Inches(1.0), Inches(6.6))

    txb(s, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.55),
        "The encoding carries no structural knowledge of the cube group - all 324 dimensions are "
        "treated as independent by the MLP. Two states related by a cube symmetry get different vectors.",
        size=11, italic=True, color=GRAY)

    set_notes(s,
        "The cube state goes through two transformations before reaching the neural network:\n\n"
        "Step 1 - cp/co/ep/eo -> 54-char facelet string:\n"
        "  The to_kociemba_string() method reconstructs which colour appears on each of the 54 "
        "facelet positions (9 per face, 6 faces). The 6 fixed centres are always their canonical colour. "
        "Corners and edges contribute their 3 or 2 sticker colours, rotated by their current orientation.\n\n"
        "Step 2 - 54-char string -> 324-dim one-hot vector:\n"
        "  Each facelet character (U/R/F/D/L/B) is mapped to index 0-5. "
        "A (54, 6) binary matrix is created with a single 1 per row and flattened to shape (324,) as float32.\n\n"
        "Why facelet encoding instead of cp/co/ep/eo directly? "
        "The facelet string is also needed by the Kociemba solver interface, "
        "so reusing it avoids a separate encoding step. "
        "The network can learn whatever internal representation it needs from the raw one-hot input.\n\n"
        "Limitation: the representation does not exploit the 48-element symmetry group of the cube. "
        "Augmenting training data with all 48 symmetry-equivalent states could reduce unique scrambles needed by 48x.")


def slide_mlp(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Neural Network Architecture",
              "Multi-layer perceptron: 3 hidden layers, ReLU activations, cross-entropy training")

    section_label(s, "Layer-by-layer breakdown", Inches(0.3), Inches(1.0), w=Inches(3.5))
    arch_data = [
        ["Layer",    "Neurons", "Activation", "Parameters"],
        ["Input",      "324",   "-",           "0"],
        ["Hidden 1",   "256",   "ReLU",        "324x256 + 256 = 83,456"],
        ["Hidden 2",   "256",   "ReLU",        "256x256 + 256 = 65,792"],
        ["Hidden 3",   "128",   "ReLU",        "256x128 + 128 = 32,896"],
        ["Output",      "18",   "Softmax*",    "128x18  + 18  =  2,322"],
        ["Total",        "",    "",             "~184,466 params"],
    ]
    tbl_create(s, arch_data, Inches(0.3), Inches(1.40), Inches(7.5), row_h=Inches(0.40))

    add_pic(s, "fig_02_mlp_architecture.png", Inches(8.0), Inches(1.0), Inches(5.0))

    set_notes(s,
        "Architecture details:\n\n"
        "Input (324): one-hot encoded facelet vector. Shape (batch, 324) in PyTorch.\n\n"
        "Three hidden layers with ReLU:\n"
        "  ReLU(x) = max(0, x). Avoids vanishing gradients compared to sigmoid/tanh.\n"
        "  No dropout, no batch-norm. The dataset is large enough at training scale "
        "that overfitting is not a significant problem for this model size.\n\n"
        "Output (18): one logit per move in MOVE_LABELS order:\n"
        "  [U, U', U2, R, R', R2, F, F', F2, D, D', D2, L, L', L2, B, B', B2]\n\n"
        "Softmax converts logits to a probability distribution at inference time. "
        "During training, PyTorch's CrossEntropyLoss applies log-softmax internally for numerical stability.\n\n"
        "Loss: L = -sum_k y_k * log(p_k) where y is a one-hot label vector. "
        "Adam optimiser with lr=0.001, betas=(0.9, 0.999). Batch size 64.\n\n"
        "Total parameter count: ~184k. Intentionally small - the goal was a model "
        "trainable in under an hour on a laptop CPU, not a state-of-the-art result.\n\n"
        "RubiksMLP is defined in solvers/ai_solver.py. Weights saved with torch.save / torch.load.")


def slide_dataset(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Dataset Generation - Expert Imitation",
              "Training labels come from a BFS expert solver; each sample is (state, next_move)")

    section_label(s, "How one training sample is created",
                  Inches(0.3), Inches(1.0), w=Inches(4.5))
    formula_box(s, Inches(0.3), Inches(1.40), Inches(6.0), Inches(2.55), [
        "1.  Scramble solved cube with d random moves",
        "    scramble = generate_scramble_at_depth(d)",
        "",
        "2.  Ask BFS expert for the full solution",
        "    solution = bidirectional_BFS(cube)",
        "",
        "3.  Record first move as label",
        "    label = MOVE_TO_INDEX[ solution[0] ]",
        "",
        "4.  Encode cube state as 324-dim vector",
        "    state = encode_state(cube)   # float32",
    ], size=11)

    section_label(s, "Dataset at each depth level", Inches(6.6), Inches(1.0), w=Inches(3.5))
    stats_data = [
        ["Depth", "New samples", "Kept (20%)", "Total size"],
        ["1",  "1,000",   "0",     "1,000"],
        ["2",  "1,000",   "200",   "1,200"],
        ["3",  "1,000",   "240",   "1,240"],
        ["4",  "1,000",   "248",   "1,248"],
        ["5",  "1,000",   "250",   "1,250"],
        ["6",  "1,000",   "250",   "1,250"],
        ["7",  "1,000",   "250",   "1,250"],
        ["8",  "1,000",   "250",   "1,250"],
    ]
    tbl_create(s, stats_data, Inches(6.6), Inches(1.40), Inches(6.4), row_h=Inches(0.38))

    section_label(s, "80% replacement rule", Inches(0.3), Inches(4.15), w=Inches(3.0))
    formula_box(s, Inches(0.3), Inches(4.55), Inches(6.0), Inches(1.45), [
        "D(d+1) = NewSamples(d+1)  ∪  Sample(D(d), keep=0.20)",
        "",
        "# Prevents easy depths from dominating the loss",
    ], size=12)

    txb(s, Inches(0.3), Inches(6.15), Inches(12.7), Inches(0.5),
        "Files saved: data/datasets/dataset_depth_N.npz  (arrays: states, labels, depths)  |  "
        "BFS timeout per sample: 5 s  |  Failed samples skipped (skip rate near zero for depth <= 5)",
        size=11, italic=True, color=GRAY)

    set_notes(s,
        "The key insight: we never need to teach the model to solve the whole cube at once. "
        "We just need it to predict the correct next move given the current state.\n\n"
        "The BFS solver (bidirectional BFS, max depth 14, 30-sec timeout) generates the full optimal "
        "solution. Only the first move of that solution is used as the training label.\n\n"
        "Why only the first move and not the full sequence?\n"
        "Because the model will be evaluated step-by-step at inference time. "
        "Using only the first-move label keeps the data clean and avoids label noise from "
        "non-optimal intermediate states.\n\n"
        "The 80% replacement rule prevents catastrophic forgetting while keeping the dataset "
        "fresh for the current difficulty level. Without it, depth-1 samples would overwhelm "
        "the loss signal from harder examples.\n\n"
        "Dataset files are numpy .npz archives with three arrays:\n"
        "  states  (N, 324) float32\n"
        "  labels  (N,)     int64   (values 0-17)\n"
        "  depths  (N,)     int32   (which depth level each sample came from)\n\n"
        "Regenerate with: python main.py generate_dataset --max-depth 8 --samples 1000")


def slide_curriculum(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Curriculum Learning and Training",
              "Start easy, get harder - advance only when the model is ready")

    section_label(s, "Training loop", Inches(0.3), Inches(1.0), w=Inches(2.5))
    formula_box(s, Inches(0.3), Inches(1.40), Inches(6.0), Inches(3.65), [
        "for d in range(1, max_depth + 1):",
        "",
        "    # 1. Generate curriculum dataset",
        "    states, labels = gen_dataset(d)",
        "",
        "    # 2. Train for fixed epochs",
        "    for epoch in range(epochs_per_depth):",
        "        L = CrossEntropy(model(states), labels)",
        "        L.backward();  optimizer.step()",
        "",
        "    # 3. Gate: only advance if ready",
        "    rate = solve_rate(model, depth=d, n=100)",
        "    if rate < threshold:  break",
    ], size=11)

    section_label(s, "Key hyperparameters", Inches(6.5), Inches(1.0), w=Inches(3.0))
    hp_data = [
        ["Parameter",       "Value",    "Notes"],
        ["Optimiser",        "Adam",     "lr=0.001, betas=(0.9, 0.999)"],
        ["Loss",             "CE",       "CrossEntropyLoss (raw logits)"],
        ["Batch size",       "64",       "Shuffled each epoch"],
        ["Epochs/depth",     "50",       "Configurable via CLI"],
        ["Samples/depth",    "1,000",    "--samples flag"],
        ["Solve threshold",  "80%",      "Min rate to advance depth"],
        ["Max depth",        "5 (def)",  "Configurable, pre-trained to 8"],
    ]
    tbl_create(s, hp_data, Inches(6.5), Inches(1.40), Inches(6.5), row_h=Inches(0.40))

    add_pic(s, "real_loss_curves.png", Inches(6.5), Inches(4.85), Inches(6.5))
    txb(s, Inches(6.5), Inches(7.05), Inches(6.5), Inches(0.35),
        "Training loss across depth levels. Spike at each depth transition is expected.",
        size=10, italic=True, color=GRAY)

    set_notes(s,
        "Curriculum learning rationale:\n"
        "Training on all depths simultaneously causes the model to focus on easy examples "
        "(they are more common and the loss is easier to minimise). Starting at depth 1 and "
        "advancing only when the solve rate exceeds the threshold forces the model to properly "
        "learn each level before facing harder ones.\n\n"
        "The solve-rate gate (if rate < threshold: break) is important. "
        "If the model is stuck at 60% on depth 3, advancing to depth 4 just adds noise. "
        "In practice, depths 1-3 are mastered quickly (within 10-20 epochs). "
        "Depth 5 requires the full 50 epochs and sometimes plateaus around 75-80%.\n\n"
        "Adam optimiser: momentum (beta1=0.9) smooths noisy gradients; adaptive learning rate "
        "(beta2=0.999) adjusts per-parameter step sizes. Learning rate 0.001 was found to work "
        "well without tuning - standard choice for most MLP classification tasks.\n\n"
        "The loss spike at each depth transition is because the dataset suddenly becomes 80% "
        "harder examples. The model briefly performs worse before adapting.\n\n"
        "CLI: python main.py train --max-depth 5 --epochs 50 --samples 1000 --threshold 0.80\n"
        "Colab notebook: train_on_colab.ipynb (GPU training, ~5 minutes for depth 1-5)")


def slide_results(prs):
    s = blank(prs)
    set_bg(s)
    title_bar(s, "Evaluation Results",
              "100 test cases per depth, max 50 moves allowed, fixed random seed")

    section_label(s, "Solve rates and average solution length",
                  Inches(0.3), Inches(1.0), w=Inches(5.0))
    res_data = [
        ["Depth", "Greedy rate", "Greedy avg moves", "Beam (w=5) rate", "Beam avg moves", "Beam avg time"],
        ["1",   "100%",  "1.0",  "100%",  "1.0",   "0.5 ms"],
        ["2",   "100%",  "2.0",  "100%",  "2.0",   "1.0 ms"],
        ["3",   "100%",  "3.0",  "100%",  "3.0",   "4.4 ms"],
        ["4",    "97%",  "4.0",   "98%",  "4.0",   "7.2 ms"],
        ["5",    "77%",  "5.1",   "89%",  "5.1",  "11.1 ms"],
        ["6",    "31%",  "6.1",   "50%",  "6.9",  "17.9 ms"],
        ["7",    "13%",  "7.0",   "23%",  "9.6",  "30.2 ms"],
        ["8",     "3%",  "8.7",    "9%",  "9.6",  "25.7 ms"],
    ]
    tbl_create(s, res_data, Inches(0.3), Inches(1.40), Inches(12.7), row_h=Inches(0.42))

    add_pic(s, "real_solve_rate_bars.png", Inches(0.3), Inches(5.0), Inches(6.3))

    takeaways = [
        (Inches(6.8),  Inches(5.0), GREEN,
         "Both solvers: perfect\nat depths 1-3"),
        (Inches(9.0),  Inches(5.0), STEEL,
         "Beam +12pp over\ngreedy at depth 5"),
        (Inches(11.1), Inches(5.0), RED,
         "Sharp drop beyond\ndepth 6 (distribution shift)"),
    ]
    for lx, ty, col, txt in takeaways:
        add_rect(s, lx, ty, Inches(2.0), Inches(1.1), STRIP,
                 RGBColor(0x99, 0xAA, 0xCC))
        txb(s, lx + Inches(0.08), ty + Inches(0.08), Inches(1.84), Inches(0.95),
            txt, size=12, bold=True, color=col)

    set_notes(s,
        "Evaluation setup:\n"
        "  - 100 scrambles per depth, seed fixed for reproducibility\n"
        "  - Scramble: exactly d moves, no consecutive same-face moves\n"
        "  - Solver gets max 50 moves to find solution\n"
        "  - Success = cube reaches is_solved() == True within 50 moves\n\n"
        "Key findings:\n\n"
        "1. Both solvers work perfectly at depths 1-3. The model memorises these effectively "
        "because there are only a few hundred distinct states at these depths.\n\n"
        "2. Beam search gives +12 percentage points at depth 5 (77% -> 89%). "
        "The improvement is largest where the greedy policy is nearly but not quite reliable.\n\n"
        "3. Sharp degradation beyond depth 6. The model was trained on depths 1-8, but with only "
        "1,000 samples per depth. At depth 6+, the state space coverage is extremely sparse. "
        "The model has not seen most of these configurations.\n\n"
        "4. Beam search average move count at depth 7 is 9.6 (vs scramble depth 7). "
        "When beam search does succeed, it takes a slightly longer path - "
        "it is not finding optimal solutions, just any solution.\n\n"
        "5. Greedy inference time: 0.4-3.4 ms. Beam search: 0.5-55 ms. "
        "Both fast enough for interactive use.")


# ---------------------------------------------------------------------------

def main():
    prs = Presentation(TEMPLATE)
    _update_front_page_date(prs)
    _remove_slides_after_first(prs)

    slide_system_overview(prs)
    slide_cube_engine(prs)
    slide_move_tables(prs)
    slide_encoding(prs)
    slide_mlp(prs)
    slide_dataset(prs)
    slide_curriculum(prs)
    slide_results(prs)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
